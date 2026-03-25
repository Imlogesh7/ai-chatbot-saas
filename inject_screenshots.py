"""
Inject real screenshots into Slide 5 (Demo slide) of the ContextIQ pitch deck.
Replaces the placeholder grey boxes with actual product screenshots.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, copy

PPTX_PATH = "/Users/ctr-lelanchelian/Desktop/Saas_product_2/ContextIQ_Investor_Pitch_Deck.pptx"
SS_DIR    = "/Users/ctr-lelanchelian/Desktop/Saas_product_2/screenshots"

DARK_BG   = RGBColor(0x0D, 0x0F, 0x1A)
CARD_BG   = RGBColor(0x12, 0x1C, 0x35)
MID_BLUE  = RGBColor(0x1A, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x24, 0x6B, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xBE, 0xD4)
SUBTEXT   = RGBColor(0x7A, 0x8E, 0xA8)

prs = Presentation(PPTX_PATH)
BLANK = prs.slide_layouts[6]

# ── Target slide is index 4 (Slide 5, 0-indexed) ──────────────────────────────
demo_slide = prs.slides[4]

# Remove all existing shapes on that slide cleanly
sp_tree = demo_slide.shapes._spTree
for sp in list(sp_tree):
    sp_tree.remove(sp)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

def add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def accent_bar(slide, color=ACCENT_BLUE, height=0.04):
    add_rect(slide, 0, 0, 13.33, height, color)

def section_tag(slide, label, x=0.35, y=0.12):
    tag = slide.shapes.add_shape(1, Inches(x), Inches(y),
                                  Inches(len(label)*0.095+0.25), Inches(0.26))
    tag.fill.solid()
    tag.fill.fore_color.rgb = ACCENT_BLUE
    tag.line.fill.background()
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label.upper()
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

def slide_number_badge(slide, num):
    add_rect(slide, 12.5, 7.1, 0.6, 0.28, MID_BLUE)
    add_text_box(slide, str(num), 12.5, 7.1, 0.6, 0.28,
                 font_size=9, color=SUBTEXT, align=PP_ALIGN.CENTER)

def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text

# ── Rebuild Slide 5 with real screenshots ────────────────────────────────────
s = demo_slide

# Full dark background
add_rect(s, 0, 0, 13.33, 7.5, DARK_BG)

# Top accent bar
accent_bar(s, ACCENT_BLUE, 0.055)

# Section tag + header
section_tag(s, "PRODUCT DEMO")
add_text_box(s, "Product Demo — Live Application Screenshots",
             0.35, 0.38, 12.5, 0.65, font_size=30, bold=True, color=WHITE)
add_text_box(s,
    "Real screenshots from the live product at frontend-ten-hazel-61.vercel.app",
    0.35, 1.1, 12.5, 0.38, font_size=14, color=ACCENT_CYAN)

# Screenshot grid: 2×2 layout
screens = [
    ("1_dashboard.png",  "Dashboard",
     "Chatbot overview, live stats,\nquick-create button"),
    ("2_documents.png",  "Document Upload",
     "PDF drag-drop + URL submit,\nreal-time processing badges"),
    ("3_chat.png",       "AI Chat Interface",
     "Live RAG-powered conversation,\nauto-scroll, message history"),
    ("4_embed.png",      "Embed Code",
     "One-click copy script tag,\ncustomizable widget color"),
]

for i, (fname, title, desc) in enumerate(screens):
    col = i % 2
    row = i // 2
    xi = 0.35 + col * 6.5
    yi = 1.62 + row * 2.88

    img_path = os.path.join(SS_DIR, fname)

    # Card background
    add_rect(s, xi, yi, 6.2, 2.72, CARD_BG)

    # Top bar with title
    bar = s.shapes.add_shape(1, Inches(xi), Inches(yi), Inches(6.2), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

    add_text_box(s, title, xi+0.15, yi+0.1, 3.5, 0.35,
                 font_size=13, bold=True, color=WHITE)
    add_text_box(s, desc, xi+0.15, yi+0.48, 2.8, 0.55,
                 font_size=10.5, color=LIGHT_GRAY)

    # Insert screenshot image on the right side of the card
    if os.path.exists(img_path):
        try:
            pic = s.shapes.add_picture(
                img_path,
                Inches(xi + 3.1),
                Inches(yi + 0.08),
                Inches(2.95),
                Inches(2.55)
            )
            print(f"  ✓ Inserted {fname}")
        except Exception as e:
            print(f"  ✗ Could not insert {fname}: {e}")
    else:
        print(f"  ✗ File not found: {img_path}")

slide_number_badge(s, 5)
add_speaker_notes(s, """SPEAKER NOTES — DEMO SLIDE (Real Screenshots)
These are live screenshots from the deployed product at frontend-ten-hazel-61.vercel.app.

Walk through each quadrant:
1. Dashboard: show the chatbot count stats, clean sidebar navigation
2. Document Upload: show processing status badges — this is where the AI pipeline starts
3. Chat: show a real conversation — emphasize the answer quality and speed (<2 seconds)
4. Embed: show how simple the integration is — one script tag, fully configured

Pro tip: open the live product on a second screen during the presentation for a live demo.""")

# Save updated file
prs.save(PPTX_PATH)
print(f"\n✅  Presentation updated with real screenshots: {PPTX_PATH}")
